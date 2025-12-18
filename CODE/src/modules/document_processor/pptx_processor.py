"""
PowerPoint Processor - Process .pptx and .ppt files.
"""

import time
from datetime import datetime
from pathlib import Path
from typing import Dict, List

from .base import (
    BaseProcessor,
    ProcessedDocument,
    DocumentMetadata,
    TextChunk,
    ImageInfo,
)


class PowerPointProcessor(BaseProcessor):
    """
    Processor for PowerPoint files (.pptx, .ppt).

    Features:
    - Extract text from slides
    - Extract speaker notes
    - Extract text from shapes and tables
    - Track slide numbers
    """

    VERSION = "1.0"

    def __init__(self, config: Dict = None):
        """
        Initialize PowerPoint Processor.

        Args:
            config: Configuration dict with options:
                - include_notes: Include speaker notes (default: True)
                - include_hidden: Include hidden slides (default: False)
        """
        super().__init__(config)
        self.include_notes = self.config.get("include_notes", True)
        self.include_hidden = self.config.get("include_hidden", False)

    def supported_extensions(self) -> List[str]:
        return [".pptx", ".ppt"]

    def process(self, file_path: str) -> ProcessedDocument:
        """Process PowerPoint file"""
        start_time = time.time()

        try:
            self.validate(file_path)
        except Exception as e:
            return self._create_error_result(file_path, e, start_time)

        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            return self._create_error_result(
                file_path,
                ImportError("python-pptx not installed. Run: pip install python-pptx"),
                start_time
            )

        try:
            prs = Presentation(file_path)

            all_text = []
            chunks = []
            images = []
            slide_count = 0

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_count += 1
                slide_text_parts = []

                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text_parts.append(shape.text.strip())

                    # Handle tables
                    if shape.has_table:
                        table_text = self._extract_table_text(shape.table)
                        if table_text:
                            slide_text_parts.append(table_text)

                # Extract speaker notes
                if self.include_notes and slide.has_notes_slide:
                    notes_frame = slide.notes_slide.notes_text_frame
                    if notes_frame and notes_frame.text.strip():
                        slide_text_parts.append(f"[Speaker Notes] {notes_frame.text.strip()}")

                # Combine slide content
                if slide_text_parts:
                    slide_text = f"\n\n--- Slide {slide_num} ---\n\n" + "\n\n".join(slide_text_parts)
                    all_text.append(slide_text)

                    # Create chunk for this slide
                    chunks.append(TextChunk(
                        text="\n\n".join(slide_text_parts),
                        start_char=0,
                        end_char=len("\n\n".join(slide_text_parts)),
                        page_number=slide_num,
                        chunk_index=slide_num - 1,
                        metadata={"slide_number": slide_num}
                    ))

            # Combine all content
            content = "\n".join(all_text)

            # Extract metadata
            metadata = self.extract_metadata(file_path)
            metadata.page_count = slide_count

            # Try to get presentation properties
            try:
                core_props = prs.core_properties
                if core_props.title:
                    metadata.title = core_props.title
                if core_props.author:
                    metadata.author = core_props.author
                if core_props.created:
                    metadata.created_date = core_props.created
                if core_props.modified:
                    metadata.modified_date = core_props.modified
            except Exception:
                pass

            metadata.extra = {
                "slide_count": slide_count,
                "word_count": len(content.split()),
                "char_count": len(content),
            }

            return ProcessedDocument(
                content=content,
                chunks=chunks,
                metadata=metadata,
                source_file=file_path,
                file_type=Path(file_path).suffix.lower().strip('.'),
                images=images,
                processed_at=datetime.now(),
                processing_time=time.time() - start_time,
                processor_version=self.VERSION,
                success=True
            )

        except Exception as e:
            return self._create_error_result(file_path, e, start_time)

    def _extract_table_text(self, table) -> str:
        """Extract text from a table shape"""
        rows_text = []
        for row in table.rows:
            cells_text = []
            for cell in row.cells:
                if cell.text.strip():
                    cells_text.append(cell.text.strip())
            if cells_text:
                rows_text.append(" | ".join(cells_text))
        return "\n".join(rows_text)
